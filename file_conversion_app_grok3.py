import gradio as gr
import os
import json
import pandas as pd
from PyPDF2 import PdfReader
from markdown2 import markdown
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from langchain.document_loaders import UnstructuredFileLoader

# Allowed file extensions
ALLOWED_EXTENSIONS = {
    '.txt', '.md', '.json', '.csv', '.pdf',
    '.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx'
}

def validate_file(filename):
    """
    Validates the file extension and checks if the file exists.
    
    Args:
        filename (str): Path to the input file.
    
    Returns:
        tuple: (bool, str) - (is_valid, error_message)
    """
    # Check if file exists
    if not os.path.exists(filename):
        return False, "Error: File does not exist."
    
    # Get file extension
    file_ext = os.path.splitext(filename)[1].lower()
    
    # Check if extension is allowed
    if file_ext not in ALLOWED_EXTENSIONS:
        return False, f"Error: Unsupported file type '{file_ext}'. Allowed types: {', '.join(ALLOWED_EXTENSIONS)}"
    
    return True, ""

def extract_text_from_txt(filename):
    """Extracts text from a plain text file."""
    with open(filename, 'r', encoding='utf-8') as f:
        return f.read()

def extract_text_from_md(filename):
    """Extracts text from a Markdown file."""
    with open(filename, 'r', encoding='utf-8') as f:
        md_content = f.read()
        return markdown(md_content, extras=['fenced-code-blocks'])

def extract_text_from_json(filename):
    """Extracts text from a JSON file."""
    with open(filename, 'r', encoding='utf-8') as f:
        json_content = json.load(f)
        return json.dumps(json_content, indent=2)

def extract_text_from_csv(filename):
    """Extracts text from a CSV file."""
    df = pd.read_csv(filename)
    return df.to_string()

def extract_text_from_pdf(filename):
    """Extracts text from a PDF file."""
    reader = PdfReader(filename)
    text = ""
    for page in reader.pages:
        text += page.extract_text() + "\n"
    return text

def extract_text_from_docx(filename):
    """Extracts text from a DOCX file."""
    doc = Document(filename)
    text = ""
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

def extract_text_from_xlsx(filename):
    """Extracts text from an XLSX file."""
    wb = load_workbook(filename)
    text = ""
    for sheet in wb:
        for row in sheet.rows:
            for cell in row:
                if cell.value:
                    text += str(cell.value) + " "
            text += "\n"
    return text

def extract_text_from_pptx(filename):
    """Extracts text from a PPTX file."""
    prs = Presentation(filename)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_with_langchain(filename):
    """Fallback: Extracts text using LangChain for unsupported file types."""
    loader = UnstructuredFileLoader(filename)
    documents = loader.load()
    return "\n".join([doc.page_content for doc in documents])

def extract_text(filename):
    """
    Extracts text from the given file based on its extension.
    
    Args:
        filename (str): Path to the input file.
    
    Returns:
        tuple: (str, str) - (extracted_text, error_message)
    """
    file_ext = os.path.splitext(filename)[1].lower()
    
    try:
        if file_ext == '.txt':
            return extract_text_from_txt(filename), ""
        elif file_ext == '.md':
            return extract_text_from_md(filename), ""
        elif file_ext == '.json':
            return extract_text_from_json(filename), ""
        elif file_ext == '.csv':
            return extract_text_from_csv(filename), ""
        elif file_ext == '.pdf':
            return extract_text_from_pdf(filename), ""
        elif file_ext == '.docx' or file_ext == '.doc':
            return extract_text_from_docx(filename), ""
        elif file_ext == '.xlsx' or file_ext == '.xls':
            return extract_text_from_xlsx(filename), ""
        elif file_ext == '.pptx' or file_ext == '.ppt':
            return extract_text_from_pptx(filename), ""
        else:
            # Fallback to LangChain for unsupported types
            return extract_text_with_langchain(filename), ""
    except Exception as e:
        return "", f"Error extracting text: {str(e)}"

def save_extracted_text(filename, text):
    """
    Saves the extracted text to a new file.
    
    Args:
        filename (str): Original filename.
        text (str): Extracted text.
    
    Returns:
        str: Success or error message.
    """
    try:
        output_filename = os.path.splitext(filename)[0] + "_extracted.txt"
        with open(output_filename, 'w', encoding='utf-8') as f:
            f.write(text)
        return f"Text successfully saved to {output_filename}"
    except Exception as e:
        return f"Error saving file: {str(e)}"

def process_file(filename):
    """
    Main function to process the input file and extract text.
    
    Args:
        filename (str): Path to the input file.
    
    Returns:
        str: Result message to display in the UI.
    """
    # Validate the file
    is_valid, error_message = validate_file(filename)
    if not is_valid:
        return error_message
    
    # Extract text from the file
    extracted_text, extraction_error = extract_text(filename)
    if extraction_error:
        return extraction_error
    
    # Save the extracted text
    save_result = save_extracted_text(filename, extracted_text)
    return save_result

def create_ui():
    """
    Creates the Gradio UI for the application.
    
    Returns:
        gr.Blocks: The Gradio interface.
    """
    with gr.Blocks() as demo:
        gr.Markdown("# File Text Extractor")
        gr.Markdown("Enter the path to a file to extract its text content.")
        
        filename_input = gr.Textbox(label="Input Filename", placeholder="e.g., /path/to/file.pdf")
        output_text = gr.Textbox(label="Result", interactive=False)
        
        submit_button = gr.Button("Extract Text")
        submit_button.click(
            fn=process_file,
            inputs=filename_input,
            outputs=output_text
        )
    
    return demo

# Launch the UI
if __name__ == "__main__":
    ui = create_ui()
    ui.launch()
